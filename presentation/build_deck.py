#!/usr/bin/env python3
"""Сборка питч-презентации JMLC под 5-минутное выступление.

Запуск:
    ./venv/bin/pip install python-pptx
    ./venv/bin/python presentation/build_deck.py

Результат: presentation/JMLC_pitch.pptx (+ заметки докладчика с таймингом).
Все цифры — из docs/METRICS.md и results/ исходного проекта.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

OUT = Path(__file__).resolve().parent / "JMLC_pitch.pptx"

# 16:9
SW, SH = Emu(12192000), Emu(6858000)
MARGIN = Emu(700000)

# Палитра
NAVY = RGBColor(0x0F, 0x1E, 0x3D)
INK = RGBColor(0x1B, 0x25, 0x38)
TEAL = RGBColor(0x14, 0xA0, 0x8F)
ORANGE = RGBColor(0xE8, 0x7A, 0x22)
GREY = RGBColor(0x5B, 0x66, 0x72)
LIGHT = RGBColor(0xF3, 0xF5, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"


def _text(tf, runs, size, color, bold=False, align=PP_ALIGN.LEFT, space_after=6):
    """runs: str | list[(str, dict)]."""
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(space_after)
    items = runs if isinstance(runs, list) else [(runs, {})]
    for txt, opt in items:
        r = p.add_run()
        r.text = txt
        f = r.font
        f.name = FONT
        f.size = Pt(opt.get("size", size))
        f.bold = opt.get("bold", bold)
        f.color.rgb = opt.get("color", color)


def _box(slide, x, y, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    return sh


def _tb(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.vertical_anchor = anchor
    return tb.text_frame


def _bullets(tf, items, size=18, color=INK, gap=10):
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        parts = it if isinstance(it, list) else [(it, {})]
        for txt, opt in parts:
            r = p.add_run()
            r.text = txt
            f = r.font
            f.name = FONT
            f.size = Pt(opt.get("size", size))
            f.bold = opt.get("bold", False)
            f.color.rgb = opt.get("color", color)


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _header(slide, kicker, title):
    _box(slide, 0, 0, SW, Emu(1500000), fill=NAVY)
    _box(slide, MARGIN, Emu(560000), Emu(180000), Emu(560000), fill=ORANGE)
    tf = _tb(slide, Emu(1000000), Emu(360000), SW - Emu(1700000), Emu(900000),
             anchor=MSO_ANCHOR.MIDDLE)
    _text(tf, kicker.upper(), 13, TEAL, bold=True, space_after=2)
    p = tf.add_paragraph()
    r = p.add_run(); r.text = title
    r.font.name = FONT; r.font.size = Pt(30); r.font.bold = True
    r.font.color.rgb = WHITE


TOTAL = 11


def _page(slide, n):
    tf = _tb(slide, SW - Emu(900000), SH - Emu(500000), Emu(700000), Emu(350000))
    _text(tf, f"{n}/{TOTAL}", 11, GREY, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────────
def build() -> None:
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    # 1. Титул
    s = _blank(prs)
    _box(s, 0, 0, SW, SH, fill=NAVY)
    _box(s, MARGIN, Emu(2400000), Emu(200000), Emu(1500000), fill=ORANGE)
    tf = _tb(s, Emu(1050000), Emu(2350000), SW - Emu(1700000), Emu(2100000),
             anchor=MSO_ANCHOR.MIDDLE)
    _text(tf, "JUNIOR ML CONTEST 2026 · AI TALENT HUB", 15, TEAL, bold=True, space_after=10)
    p = tf.add_paragraph(); r = p.add_run()
    r.text = "Классификатор городских участков\nпо уличным панорамам"
    r.font.name = FONT; r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.space_before = Pt(18)
    r2 = p2.add_run(); r2.text = "Дмитрий Новичков · ПИК AI Lab"
    r2.font.name = FONT; r2.font.size = Pt(20); r2.font.color.rgb = RGBColor(0xC7, 0xD0, 0xDC)
    _notes(s, """[0:00–0:15] Здравствуйте. Меня зовут Дмитрий Новичков. Я покажу проект,
который по уличной панораме определяет тип городского участка. Это боль, с которой
я сталкиваюсь в работе аналитика девелопера. Пять минут — поехали.""")

    # 2. Проблема
    s = _blank(prs)
    _box(s, 0, 0, SW, SH, fill=WHITE)
    _header(s, "Проблема", "Тысячи участков нужно оценивать вручную")
    tf = _tb(s, MARGIN, Emu(1900000), SW - 2 * MARGIN, Emu(4200000))
    _bullets(tf, [
        [("Аналитики руками просматривают тысячи участков Москвы: стройка, "
          "деградация, недоиспользование. Медленно и не масштабируется.", {})],
        [("Спутник не помогает — ", {}), ("не видит дворы и фасады.", {"bold": True})],
        [("Текст в GPKG часто не различает класс: ", {}),
         ("1712 участков", {"bold": True, "color": ORANGE}),
         (" с одной фразой «зелень во дворах» — а классы разные.", {})],
        [("Отличить их можно только по виду с улицы.", {"bold": True, "color": TEAL})],
    ], size=20, gap=16)
    _page(s, 2)
    _notes(s, """[0:15–0:45] Проблема простая и реальная. Аналитики вручную относят
тысячи участков к типам застройки. Спутник тут бесполезен — он не видит дворы и
фасады. Текстовое описание из кадастра тоже подводит: например, 1712 участков имеют
одинаковую фразу «зелень во дворах», но по факту это разные классы. Различить их
можно только по уличной панораме. Отсюда идея проекта.""")

    # 3. Решение + классы
    s = _blank(prs)
    _box(s, 0, 0, SW, SH, fill=WHITE)
    _header(s, "Решение", "Панорама участка → тип территории")
    tf = _tb(s, MARGIN, Emu(1750000), SW - 2 * MARGIN, Emu(700000))
    _text(tf, [("GeoPackage → панорамы Яндекса → CLIP-классификатор → ", {"size": 19}),
               ("6 классов UTT", {"size": 19, "bold": True, "color": TEAL}),
               (" + веб-инструмент для правок.", {"size": 19})], 19, INK)
    classes = [
        ("active_urban", "сформированная застройка"),
        ("active_construction", "активная стройка"),
        ("frozen_construction", "замороженная стройка"),
        ("low_density_degraded", "малоэтажная / деградирующая"),
        ("underused_infrastructure", "недоиспользуемая инфра"),
        ("natural_areas", "природные территории"),
    ]
    cw = (SW - 2 * MARGIN - Emu(400000)) // 3
    ch = Emu(1250000)
    x0, y0 = MARGIN, Emu(2700000)
    for i, (code, desc) in enumerate(classes):
        cx = x0 + (i % 3) * (cw + Emu(200000))
        cy = y0 + (i // 3) * (ch + Emu(250000))
        _box(s, cx, cy, cw, ch, fill=LIGHT, line=RGBColor(0xDD, 0xE3, 0xEA))
        _box(s, cx, cy, Emu(120000), ch, fill=TEAL)
        t = _tb(s, cx + Emu(240000), cy + Emu(180000), cw - Emu(380000),
                ch - Emu(360000), anchor=MSO_ANCHOR.MIDDLE)
        _text(t, code, 15, NAVY, bold=True, space_after=4)
        p = t.add_paragraph(); r = p.add_run(); r.text = desc
        r.font.name = FONT; r.font.size = Pt(13); r.font.color.rgb = GREY
    _page(s, 3)
    _notes(s, """[0:45–1:15] Решение: берём координаты участка из GeoPackage, находим
уличную панораму, вырезаем кадр и предсказываем один из шести типов территории —
от активной застройки до природных зон. Вокруг модели — веб-инструмент, где аналитик
проверяет и правит результат. То есть это не просто модель, а рабочий продукт.""")

    # 4. Данные и качество
    s = _blank(prs)
    _box(s, 0, 0, SW, SH, fill=WHITE)
    _header(s, "Данные и EDA", "Главный рычаг качества — чистота меток")
    tf = _tb(s, MARGIN, Emu(1850000), (SW - 2 * MARGIN) * 58 // 100, Emu(4000000))
    _bullets(tf, [
        [("29% кропов", {"bold": True, "color": ORANGE}),
         (" смотрели «мимо» участка (bearing≈0) — чиним мульти-азимутом 0/90/180/270.", {})],
        [("Согласие OSM-меток с содержимым кадра — ", {}),
         ("~20%", {"bold": True, "color": ORANGE}),
         (". Переразметили VLM-судьёй (Qwen) + ручная вычитка.", {})],
        [("Чистка: дубли, размытые, нерелевантные ракурсы (perceptual hash + CLIP).", {})],
        [("Сплит по ", {}), ("object_id", {"bold": True}),
         (" — без утечки сезонов между train и val.", {})],
    ], size=18, gap=14)
    # right stat cards
    stats = [("7535", "участков GPKG\n(валидация продукта)"),
             ("~600", "объектов собрано\nпо 6 классам"),
             ("~2.4k", "кропов после\nмульти-хединга и чистки")]
    rx = MARGIN + (SW - 2 * MARGIN) * 62 // 100
    rw = (SW - 2 * MARGIN) * 38 // 100
    for i, (num, lab) in enumerate(stats):
        cy = Emu(1850000) + i * Emu(1150000)
        _box(s, rx, cy, rw, Emu(1000000), fill=NAVY)
        t = _tb(s, rx + Emu(260000), cy + Emu(120000), rw - Emu(400000), Emu(760000),
                anchor=MSO_ANCHOR.MIDDLE)
        _text(t, num, 30, ORANGE, bold=True, space_after=2)
        p = t.add_paragraph(); r = p.add_run(); r.text = lab
        r.font.name = FONT; r.font.size = Pt(12); r.font.color.rgb = WHITE
    _page(s, 4)
    _notes(s, """[1:15–1:55] Про данные — здесь была основная работа. Первое: почти
треть кадров смотрели мимо участка, потому что панорама снята прямо в точке объекта.
Починили, рендеря четыре направления. Второе и главное: сырые метки OSM совпадали с
тем, что реально на кадре, лишь на 20%. Поэтому переразметили данные VLM-судьёй и
вычитали руками. Плюс дедупликация и валидация без утечки сезонов. Именно качество
меток, а не архитектура, дало основной прирост.""")

    # 5. VLM — двойная роль
    s = _blank(prs)
    _box(s, 0, 0, SW, SH, fill=WHITE)
    _header(s, "Роль VLM", "Разметка данных и объяснимость")
    colw = (SW - 2 * MARGIN - Emu(300000)) // 2
    cy = Emu(1850000)
    chh = Emu(3600000)
    # left: weak supervision
    _box(s, MARGIN, cy, colw, chh, fill=LIGHT, line=RGBColor(0xDD, 0xE3, 0xEA))
    _box(s, MARGIN, cy, colw, Emu(150000), fill=TEAL)
    t = _tb(s, MARGIN + Emu(260000), cy + Emu(300000), colw - Emu(500000), chh - Emu(450000))
    _text(t, "1 · VLM-судья для разметки", 18, NAVY, bold=True, space_after=12)
    _bullets(t, [
        [("OSM-метки ненадёжны (согласие с кадром ", {}),
         ("~20%", {"bold": True, "color": ORANGE}), (")", {})],
        [("Qwen-VL переразмечает кроп по содержимому — weak supervision", {})],
        [("Спорные случаи — ручная вычитка", {})],
        [("Итог: accuracy ", {}), ("0.27 → 0.76", {"bold": True, "color": TEAL})],
    ], size=15, gap=11)
    # right: generative explainability
    rx = MARGIN + colw + Emu(300000)
    _box(s, rx, cy, colw, chh, fill=LIGHT, line=RGBColor(0xDD, 0xE3, 0xEA))
    _box(s, rx, cy, colw, Emu(150000), fill=ORANGE)
    t2 = _tb(s, rx + Emu(260000), cy + Emu(300000), colw - Emu(500000), chh - Emu(450000))
    _text(t2, "2 · Генеративный VLM (объяснимость)", 18, NAVY, bold=True, space_after=12)
    _bullets(t2, [
        [("PaliGemma 3B + LoRA: класс ", {}),
         ("+ bbox-доказательство", {"bold": True, "color": ORANGE}),
         (" на кадре", {})],
        [("Evidence IoU ", {}), ("0.26", {"bold": True}),
         (" — показывает, «почему» такой класс", {})],
        [("Как классификатор слабее probe (acc 0.21) — честно", {})],
        [("Ценность — интерпретируемость для аналитика", {})],
    ], size=15, gap=11)
    tf = _tb(s, MARGIN, Emu(5550000), SW - 2 * MARGIN, Emu(700000))
    _text(tf, [("Вывод: VLM отвечает за качество меток и объяснимость, "
                "а продакшен-классификатор — лёгкий CLIP-probe.",
                {"size": 16, "color": GREY})], 16, GREY)
    _page(s, 5)
    _notes(s, """[1:55–2:30] Отдельно про VLM — у него в проекте две роли. Первая:
VLM-судья на базе Qwen переразмечает кадры по их реальному содержимому вместо шумных
OSM-меток — это weak supervision, спорное я вычитывал руками. Именно это подняло
точность с 0.27 до 0.76. Вторая роль — объяснимость: дообученная LoRA поверх
PaliGemma выдаёт не только класс, но и рамку-доказательство на кадре, почему так.
Как самостоятельный классификатор генеративный VLM слабее probe, и я это честно
показываю; его ценность — интерпретируемость. В проде классифицирует лёгкий
CLIP-probe, а VLM работает на качество данных и объяснения.""")

    # 6. Как работает
    s = _blank(prs)
    _box(s, 0, 0, SW, SH, fill=WHITE)
    _header(s, "Как работает", "Замороженный CLIP + линейный пробинг")
    steps = ["Контур участка\n(GeoPackage)", "Точки панорам\nURL / грани полигона",
             "Кроп 896×672\n4 азимута", "CLIP ViT-L-14\n(заморожен)",
             "Linear probe\n6 классов", "Голосование\nпо ракурсам"]
    n = len(steps)
    gap = Emu(150000)
    bw = (SW - 2 * MARGIN - gap * (n - 1)) // n
    by = Emu(2550000)
    for i, st in enumerate(steps):
        bx = MARGIN + i * (bw + gap)
        col = TEAL if i in (3, 4) else NAVY
        _box(s, bx, by, bw, Emu(1300000), fill=col)
        t = _tb(s, bx + Emu(80000), by, bw - Emu(160000), Emu(1300000),
                anchor=MSO_ANCHOR.MIDDLE)
        _text(t, st, 13, WHITE, bold=True, align=PP_ALIGN.CENTER)
        if i < n - 1:
            a = _tb(s, bx + bw - Emu(40000), by + Emu(480000), Emu(220000), Emu(350000),
                    anchor=MSO_ANCHOR.MIDDLE)
            _text(a, "›", 26, ORANGE, bold=True, align=PP_ALIGN.CENTER)
    tf = _tb(s, MARGIN, Emu(4400000), SW - 2 * MARGIN, Emu(900000))
    _text(tf, [("Пробинг поверх замороженного CLIP устойчив к переобучению на малых "
                "данных; у участка несколько ракурсов — класс выбирается взвешенным "
                "голосованием по уверенности.", {"size": 16, "color": GREY})], 16, GREY)
    _page(s, 6)
    _notes(s, """[2:30–3:00] Как это работает под капотом. По контуру участка находим
точки панорам, вырезаем кадры в четырёх направлениях, прогоняем через замороженный
CLIP и обучаем поверх него лёгкий линейный классификатор. Заморозка важна: данных
немного, а так модель не переобучается. У участка обычно несколько ракурсов —
финальный класс выбираем голосованием с весом по уверенности.""")

    # 6. Результаты
    s = _blank(prs)
    _box(s, 0, 0, SW, SH, fill=WHITE)
    _header(s, "Результаты", "Продакшен-модель на честной валидации")
    # big metrics
    metrics = [("76.2%", "Accuracy"), ("61.0%", "macro-F1"), ("73.8%", "balanced F1")]
    mw = (SW - 2 * MARGIN - Emu(400000)) // 3
    for i, (num, lab) in enumerate(metrics):
        mx = MARGIN + i * (mw + Emu(200000))
        _box(s, mx, Emu(1850000), mw, Emu(1350000), fill=LIGHT,
             line=RGBColor(0xDD, 0xE3, 0xEA))
        t = _tb(s, mx, Emu(1850000), mw, Emu(1350000), anchor=MSO_ANCHOR.MIDDLE)
        _text(t, num, 40, TEAL, bold=True, align=PP_ALIGN.CENTER, space_after=2)
        p = t.add_paragraph(); p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = lab
        r.font.name = FONT; r.font.size = Pt(15); r.font.color.rgb = GREY
    tf = _tb(s, MARGIN, Emu(3450000), SW - 2 * MARGIN, Emu(500000))
    _text(tf, [("CLIP ViT-L-14 probe · n_val = 239 · сплит по object_id · все 6 классов",
                {"size": 14, "color": GREY})], 14, GREY)
    # evolution
    tf2 = _tb(s, MARGIN, Emu(4050000), SW - 2 * MARGIN, Emu(400000))
    _text(tf2, "Путь к результату (рычаг — метки, не архитектура):", 16, NAVY, bold=True)
    evo = ["zero-shot 0.20", "fine-tune 0.40", "OSM-probe 0.27",
           "VLM-probe 0.76", "3 класса 0.82"]
    ew = (SW - 2 * MARGIN - Emu(320000)) // 5
    ey = Emu(4600000)
    for i, e in enumerate(evo):
        ex = MARGIN + i * (ew + Emu(80000))
        hi = "0.76" in e
        _box(s, ex, ey, ew, Emu(720000), fill=ORANGE if hi else NAVY)
        t = _tb(s, ex + Emu(60000), ey, ew - Emu(120000), Emu(720000),
                anchor=MSO_ANCHOR.MIDDLE)
        _text(t, e, 13, WHITE, bold=hi, align=PP_ALIGN.CENTER)
    _page(s, 7)
    _notes(s, """[3:00–3:40] Результаты. Продакшен-модель на честной валидации со
сплитом по объектам: точность 76%, macro-F1 61%, а сбалансированный F1 — почти 74%.
В валидации есть все шесть классов. Посмотрите на нижнюю ленту: мы начинали с
zero-shot CLIP на уровне 20%, fine-tune на шумных метках упирался в 40%, а переход на
VLM-метки поднял точность до 76%. То есть главный рычаг — качество данных.""")

    # 7. Честность / ограничения
    s = _blank(prs)
    _box(s, 0, 0, SW, SH, fill=WHITE)
    _header(s, "Честно об ограничениях", "Где ML сильнее текста, а где — нет")
    tf = _tb(s, MARGIN, Emu(1900000), SW - 2 * MARGIN, Emu(3800000))
    _bullets(tf, [
        [("На всём GPKG текстовая эвристика даёт ", {}),
         ("87.6%", {"bold": True, "color": ORANGE}),
         (" — выше ML. И это ожидаемо: текст кадастра уже содержит сильные подсказки.", {})],
        [("Но на кластере «зелень во дворах» текст бессилен — там класс различает "
          "только панорама. Это и есть ниша ML.", {"color": TEAL, "bold": True})],
        [("Ограничения честно задокументированы: domain shift (обучение на research-"
          "домене), малые выборки с панорамой (25–75 объектов).", {})],
        [("Вывод: ML — дополнение к тексту на трудных кластерах, а не замена эвристики.",
          {"bold": True})],
    ], size=19, gap=16)
    _page(s, 8)
    _notes(s, """[3:40–4:05] Важно быть честным. Если сравнивать в лоб на всём GPKG,
простая текстовая эвристика даёт 88% — выше нашей модели. Это нормально: в тексте
кадастра уже зашиты подсказки. Ценность ML в другом — там, где текст одинаковый, как
на «зелени во дворах», решает только картинка. Я явно документирую ограничения:
сдвиг домена и маленькие выборки. ML — это дополнение к тексту на сложных случаях.""")

    # 8. Продукт
    s = _blank(prs)
    _box(s, 0, 0, SW, SH, fill=WHITE)
    _header(s, "Продукт", "Веб-инструмент аналитика с обратной связью")
    flow = ["Загрузка GPKG", "Автокласс по панораме", "Карта + галерея ракурсов",
            "Правка аналитиком", "Экспорт GPKG"]
    for i, f in enumerate(flow):
        cy = Emu(1900000) + i * Emu(760000)
        _box(s, MARGIN, cy, Emu(560000), Emu(600000), fill=TEAL)
        t = _tb(s, MARGIN, cy, Emu(560000), Emu(600000), anchor=MSO_ANCHOR.MIDDLE)
        _text(t, str(i + 1), 22, WHITE, bold=True, align=PP_ALIGN.CENTER)
        t2 = _tb(s, MARGIN + Emu(760000), cy, (SW - 2 * MARGIN) * 55 // 100, Emu(600000),
                 anchor=MSO_ANCHOR.MIDDLE)
        _text(t2, f, 19, INK, bold=True)
    # feedback card
    rx = MARGIN + (SW - 2 * MARGIN) * 64 // 100
    rw = (SW - 2 * MARGIN) * 36 // 100
    _box(s, rx, Emu(1900000), rw, Emu(3400000), fill=NAVY)
    t = _tb(s, rx + Emu(280000), Emu(2150000), rw - Emu(500000), Emu(3000000))
    _text(t, "ПЕТЛЯ ОБРАТНОЙ СВЯЗИ", 14, ORANGE, bold=True, space_after=12)
    _bullets(t, [
        [("Правки аналитиков = gold-метки", {"color": WHITE})],
        [("7535 участков", {"color": ORANGE, "bold": True}),
         (" с корректировками", {"color": WHITE})],
        [("→ дообучение модели", {"color": WHITE})],
        [("Flask + Leaflet, экспорт в GeoPackage", {"color": RGBColor(0xC7, 0xD0, 0xDC),
                                                    "size": 14})],
    ], size=16, gap=12)
    _page(s, 9)
    _notes(s, """[4:05–4:30] Теперь продукт. Аналитик загружает GeoPackage, видит
автоклассификацию по панораме, карту и все ракурсы, правит класс и выгружает
результат обратно. Ключевое — правки не пропадают: это готовые gold-метки. Уже есть
7535 скорректированных участков, которые возвращаются в дообучение. Получается
замкнутый цикл: продукт улучшает модель.""")

    # 9. Инженерия / MLOps / агент
    s = _blank(prs)
    _box(s, 0, 0, SW, SH, fill=WHITE)
    _header(s, "Инженерия и MLOps", "Воспроизводимо и по-взрослому")
    cards = [
        ("Стек", "Python · PyTorch · OpenCLIP\nFlask · GeoPandas"),
        ("Конфиги и трекинг", "Hydra · TensorBoard\nW&B · MLflow"),
        ("Пайплайн", "DVC · Makefile\n7 модульных шагов"),
        ("Качество", "pytest · ruff · black\npre-commit"),
        ("Доставка", "Docker · GitHub Actions CI\ndemo-данные для CI"),
        ("AI-агент", "LangGraph: evaluate →\nplan → act до цели"),
    ]
    cw = (SW - 2 * MARGIN - Emu(400000)) // 3
    ch = Emu(1350000)
    for i, (h, b) in enumerate(cards):
        cx = MARGIN + (i % 3) * (cw + Emu(200000))
        cy = Emu(1900000) + (i // 3) * (ch + Emu(250000))
        _box(s, cx, cy, cw, ch, fill=LIGHT, line=RGBColor(0xDD, 0xE3, 0xEA))
        _box(s, cx, cy, cw, Emu(140000), fill=TEAL if i != 5 else ORANGE)
        t = _tb(s, cx + Emu(240000), cy + Emu(260000), cw - Emu(400000),
                ch - Emu(400000))
        _text(t, h, 16, NAVY, bold=True, space_after=6)
        p = t.add_paragraph(); r = p.add_run(); r.text = b
        r.font.name = FONT; r.font.size = Pt(13); r.font.color.rgb = GREY
    _page(s, 10)
    _notes(s, """[4:30–4:50] Инженерно всё сделано по-взрослому и воспроизводимо:
конфигурация на Hydra, трекинг экспериментов в TensorBoard, W&B и MLflow, стадии
пайплайна в DVC, тесты и линтеры на pre-commit, Docker и CI на каждый коммит.
Отдельно — автономный LangGraph-агент, который сам гоняет цикл «оцени — спланируй —
обучи» до целевых метрик.""")

    # 10. Итог
    s = _blank(prs)
    _box(s, 0, 0, SW, SH, fill=NAVY)
    _box(s, MARGIN, Emu(1700000), Emu(200000), Emu(1600000), fill=ORANGE)
    tf = _tb(s, Emu(1050000), Emu(1650000), SW - Emu(1700000), Emu(2000000))
    _text(tf, "ИТОГ", 15, TEAL, bold=True, space_after=10)
    _bullets(tf, [
        [("Полный цикл: данные → EDA → CLIP-probe → веб-пилот → AI-агент", {"color": WHITE})],
        [("76.2% на честной валидации; ниша — где текст не различает класс",
          {"color": WHITE})],
        [("Работающий продукт с петлёй обратной связи для аналитиков", {"color": WHITE})],
    ], size=19, gap=14)
    _box(s, MARGIN, Emu(4900000), SW - 2 * MARGIN, Emu(2), line=TEAL)
    tf2 = _tb(s, MARGIN, Emu(5150000), SW - 2 * MARGIN, Emu(900000))
    _text(tf2, [("Дмитрий Новичков   ", {"size": 18, "bold": True, "color": WHITE}),
                ("novichkovde@pik.ru · github.com/mfclabber/aith_junior_ml_contest",
                 {"size": 16, "color": RGBColor(0xC7, 0xD0, 0xDC)})], 18, WHITE)
    tf3 = _tb(s, MARGIN, Emu(5750000), SW - 2 * MARGIN, Emu(500000))
    _text(tf3, "Демо: make demo  ·  Спасибо! Готов к вопросам.", 15, TEAL, bold=True)
    _notes(s, """[4:50–5:00] Итог: это полный цикл от сбора данных до продукта и
автономного агента, с честными метриками и понятной нишей применения. Спасибо!
Готов ответить на вопросы. Демо поднимается одной командой.""")

    prs.save(OUT)
    print(f"saved {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
